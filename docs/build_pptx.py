"""
Build the final presentation for DS 216o: Applied AI in Healthcare.
Run from project root:  .venv/Scripts/python docs/build_pptx.py
"""

import sys, copy
sys.stdout.reconfigure(encoding="utf-8")

from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
import lxml.etree as etree

# ── Paths ──────────────────────────────────────────────────────────────────────
ROOT   = Path(__file__).resolve().parent.parent
DOCS   = ROOT / "docs"
TMPL   = DOCS / "AIMI-project_final_presentation_template.pptx"
OUT    = DOCS / "Automated Pneumonia and COVID-19 Detection from Chest X-Ray Images.pptx"
IMGS   = ROOT
SAMPLE = ROOT / "test_images"

# ── Slide canvas (widescreen 16:9) ─────────────────────────────────────────────
# 13.33" x 7.50"  — content safe zone starts at T≈1.55" (below title)
CONTENT_T  = 1.55   # top of content area
CONTENT_H  = 5.75   # height of content area
CONTENT_L  = 0.40   # left margin
CONTENT_W  = 12.50  # usable width


# ═══════════════════════════════════════════════════════════════════════════════
# ── Low-level helpers ──────────────────────────────────────────────────────────

def _remove_content_ph(slide):
    """Remove the idx=1 content placeholder from a slide (we place our own elements)."""
    for ph in list(slide.placeholders):
        if ph.placeholder_format.idx == 1:
            ph._element.getparent().remove(ph._element)
            break


def _set_title(slide, text, size=26):
    """Write text into the title placeholder (idx=0)."""
    ph = slide.placeholders[0]
    tf = ph.text_frame
    tf.clear()
    para = tf.paragraphs[0]
    run = para.add_run()
    run.text = text
    run.font.bold = True
    run.font.size = Pt(size)


def _set_cell(table, row, col, text,
              bold=False, size=12, align=PP_ALIGN.LEFT,
              color_rgb=None, bg_hex=None):
    """Set text and optional styling in a table cell."""
    cell = table.cell(row, col)
    cell.text = ""
    tf = cell.text_frame
    tf.word_wrap = True
    para = tf.paragraphs[0]
    para.alignment = align
    run = para.add_run()
    run.text = text
    run.font.bold = bold
    run.font.size = Pt(size)
    if color_rgb:
        run.font.color.rgb = RGBColor(*color_rgb)
    if bg_hex:
        tc = cell._tc
        tcPr = tc.get_or_add_tcPr()
        # remove any existing fill
        for child in list(tcPr):
            tag = etree.QName(child.tag).localname
            if tag in ("solidFill", "gradFill", "noFill", "pattFill"):
                tcPr.remove(child)
        sf = etree.SubElement(tcPr, qn("a:solidFill"))
        clr = etree.SubElement(sf, qn("a:srgbClr"))
        clr.set("val", bg_hex)


def _add_textbox(slide, left, top, width, height,
                 paragraphs, word_wrap=True):
    """
    Add a textbox at the given position (all in inches).
    paragraphs: list of dicts with keys:
        text      : str
        bold      : bool  (default False)
        size      : int   (pt, default 13)
        color     : (r,g,b) tuple or None
        indent    : int   (0=header, 1=bullet, 2=sub-bullet)
        align     : PP_ALIGN constant (default LEFT)
        space_before: int pt (default 0)
    """
    txb = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txb.text_frame
    tf.word_wrap = word_wrap

    BULLET   = "  \u2022  "   # •
    SUBBULLET = "      \u2013  "  # –

    first = True
    for p in paragraphs:
        if first:
            para = tf.paragraphs[0]
            first = False
        else:
            para = tf.add_paragraph()

        indent = p.get("indent", 1)
        text   = p.get("text", "")
        bold   = p.get("bold", False)
        size   = p.get("size", 13)
        color  = p.get("color", None)
        align  = p.get("align", PP_ALIGN.LEFT)
        space  = p.get("space_before", 0)

        para.alignment = align

        # Space before (paragraph spacing)
        if space:
            pPr = para._pPr
            if pPr is None:
                pPr = para._p.get_or_add_pPr()
            spcBef = etree.SubElement(pPr, qn("a:spcBef"))
            spcPts = etree.SubElement(spcBef, qn("a:spcPts"))
            spcPts.set("val", str(space * 100))

        if indent == 0:
            prefix = ""
        elif indent == 1:
            prefix = BULLET
        else:
            prefix = SUBBULLET

        run = para.add_run()
        run.text = prefix + text
        run.font.bold = bold
        run.font.size = Pt(size)
        if color:
            run.font.color.rgb = RGBColor(*color)

    return txb


def _add_img(slide, path, left, top, width=None, height=None):
    if not Path(path).exists():
        print(f"  [WARN] image not found: {path}")
        return
    kw = {}
    if width:  kw["width"]  = Inches(width)
    if height: kw["height"] = Inches(height)
    slide.shapes.add_picture(str(path), Inches(left), Inches(top), **kw)


def _add_caption(slide, text, left, top, width, size=10):
    """Add a small italic caption label under an image."""
    _add_textbox(slide, left, top, width, 0.25, [
        {"text": text, "bold": False, "size": size,
         "align": PP_ALIGN.CENTER, "indent": 0,
         "color": (80, 80, 80)}
    ])


# ── Color palette ──────────────────────────────────────────────────────────────
DARK_BLUE  = "1F4E79"      # table header background
MID_BLUE   = "2E75B6"      # accent (section headers)
LIGHT_BLUE = "D6E4F0"      # alternating row
HIGHLIGHT  = "E2EFDA"      # our result row (green)
WHITE      = "FFFFFF"
DARK_TEXT  = (30,  30,  30)
WHITE_TEXT = (255, 255, 255)
BLUE_TEXT  = (31,  78,  121)


# ═══════════════════════════════════════════════════════════════════════════════
prs = Presentation(str(TMPL))
slides = prs.slides


# ── SLIDE 1 ── Title ───────────────────────────────────────────────────────────
s = slides[0]
# Title placeholder (idx=0)
ph0 = s.placeholders[0]
tf = ph0.text_frame
tf.clear()
para = tf.paragraphs[0]
para.alignment = PP_ALIGN.CENTER
run = para.add_run()
run.text = "Automated Pneumonia and COVID-19 Detection\nfrom Chest X-Ray Images"
run.font.bold = True
run.font.size = Pt(32)

# Subtitle placeholder (idx=1)
ph1 = s.placeholders[1]
tf1 = ph1.text_frame
tf1.clear()
for line, sz, bld in [
    ("Rajib Roy",                          20, True),
    ("SR No: 24459  |  IISc Bengaluru",    16, False),
    ("DS 216o: Applied AI in Healthcare",  16, False),
    ("April 2026",                          14, False),
]:
    if tf1.paragraphs[-1].text == "":
        para = tf1.paragraphs[-1]
    else:
        para = tf1.add_paragraph()
    para.alignment = PP_ALIGN.CENTER
    run = para.add_run()
    run.text = line
    run.font.bold = bld
    run.font.size = Pt(sz)


# ── SLIDE 2 ── Problem Statement ──────────────────────────────────────────────
s = slides[1]
_set_title(s, "PROBLEM STATEMENT")
_remove_content_ph(s)

_add_textbox(s, CONTENT_L, CONTENT_T, CONTENT_W, CONTENT_H, [
    {"text": "Primary Objective",
     "bold": True, "size": 17, "color": BLUE_TEXT, "indent": 0},
    {"text": "3-class chest X-ray classification: Normal / Pneumonia / COVID-19",
     "bold": False, "size": 14, "indent": 1},
    {"text": "Enable automated, scalable screening to assist radiologists",
     "bold": False, "size": 14, "indent": 1},

    {"text": "Secondary Objectives", "bold": True, "size": 17,
     "color": BLUE_TEXT, "indent": 0, "space_before": 8},
    {"text": "Uncertainty quantification (MC Dropout) — auto-refer low-confidence cases",
     "bold": False, "size": 14, "indent": 1},
    {"text": "Visual explainability via Grad-CAM heat-maps on each model",
     "bold": False, "size": 14, "indent": 1},
    {"text": "Deployment as interactive Streamlit web application",
     "bold": False, "size": 14, "indent": 1},

    {"text": "Key Challenges", "bold": True, "size": 17,
     "color": BLUE_TEXT, "indent": 0, "space_before": 8},
    {"text": "Class imbalance — COVID accounts for only ~13.4% of data",
     "bold": False, "size": 14, "indent": 1},
    {"text": "Multi-source datasets requiring unified label mapping and deduplication",
     "bold": False, "size": 14, "indent": 1},
    {"text": "Clinical interpretability: model decisions must be explainable",
     "bold": False, "size": 14, "indent": 1},
    {"text": "Uncertainty must be communicated to avoid confident wrong predictions",
     "bold": False, "size": 14, "indent": 1},
])


# ── SLIDE 3 ── Methodology ────────────────────────────────────────────────────
s = slides[2]
_set_title(s, "METHODOLOGY")
_remove_content_ph(s)

_add_textbox(s, CONTENT_L, CONTENT_T, 6.0, CONTENT_H, [
    {"text": "Model Architectures (ImageNet pretrained)",
     "bold": True, "size": 16, "color": BLUE_TEXT, "indent": 0},
    {"text": "DenseNet121  —  7.2M params  —  dense skip connections",
     "bold": False, "size": 13, "indent": 1},
    {"text": "ResNet50     —  24.0M params —  residual connections",
     "bold": False, "size": 13, "indent": 1},
    {"text": "EfficientNetB0 — 4.3M params —  compound scaling",
     "bold": False, "size": 13, "indent": 1},
    {"text": "Shared head: BN\u2192Drop\u2192FC(f,256)\u2192ReLU\u2192BN\u2192Drop\u2192FC(256,3)",
     "bold": False, "size": 12, "indent": 2},

    {"text": "Two-Phase Transfer Learning",
     "bold": True, "size": 16, "color": BLUE_TEXT, "indent": 0, "space_before": 8},
    {"text": "Phase 1 (5 ep, lr=1e-3): frozen backbone — train head only",
     "bold": False, "size": 13, "indent": 1},
    {"text": "Phase 2 (\u226420 ep, lr=1e-5): unfreeze last block + joint fine-tuning",
     "bold": False, "size": 13, "indent": 1},
    {"text": "EarlyStopping (patience=5), ReduceLROnPlateau, ModelCheckpoint",
     "bold": False, "size": 12, "indent": 2},

    {"text": "Loss & Class Balancing",
     "bold": True, "size": 16, "color": BLUE_TEXT, "indent": 0, "space_before": 8},
    {"text": "Focal Loss (\u03b3=2, \u03b1=0.25) + inverse-frequency class weights",
     "bold": False, "size": 13, "indent": 1},
    {"text": "COVID weight = 2.49\u00d7  (vs Normal=0.77, Pneumonia=0.77)",
     "bold": False, "size": 12, "indent": 2},

    {"text": "Ensemble + Explainability",
     "bold": True, "size": 16, "color": BLUE_TEXT, "indent": 0, "space_before": 8},
    {"text": "AUC-weighted soft voting across all 3 models",
     "bold": False, "size": 13, "indent": 1},
    {"text": "MC Dropout (50 passes): refer if combined uncertainty > 0.35",
     "bold": False, "size": 13, "indent": 1},
    {"text": "Grad-CAM on last conv block \u2014 heat-maps overlaid on CXR",
     "bold": False, "size": 13, "indent": 1},
])

# Methodology pipeline diagram (right side)
_add_textbox(s, 6.6, CONTENT_T, 6.0, CONTENT_H, [
    {"text": "Pipeline Overview",
     "bold": True, "size": 15, "color": BLUE_TEXT, "indent": 0, "align": PP_ALIGN.CENTER},
    {"text": "",  "bold": False, "size": 6, "indent": 0},
    {"text": "\u25ba  Raw CXR images (Kermany + COVID-19 DB)",
     "bold": False, "size": 13, "indent": 0},
    {"text": "\u25ba  Unify labels \u2192 deduplicate \u2192 80/10/10 split",
     "bold": False, "size": 13, "indent": 0},
    {"text": "\u25ba  Augmentation + ImageNet normalisation",
     "bold": False, "size": 13, "indent": 0},
    {"text": "\u25ba  Phase 1: train head (backbone frozen)",
     "bold": False, "size": 13, "indent": 0},
    {"text": "\u25ba  Phase 2: fine-tune last block + head",
     "bold": False, "size": 13, "indent": 0},
    {"text": "\u25ba  Ensemble: AUC-weighted soft voting",
     "bold": False, "size": 13, "indent": 0},
    {"text": "\u25ba  MC Dropout \u2192 uncertainty \u2192 referral flag",
     "bold": False, "size": 13, "indent": 0},
    {"text": "\u25ba  Grad-CAM \u2192 saliency heat-map overlay",
     "bold": False, "size": 13, "indent": 0},
    {"text": "\u25ba  Streamlit app: single & batch inference",
     "bold": False, "size": 13, "indent": 0},
])


# ── SLIDE 4 ── Dataset ────────────────────────────────────────────────────────
s = slides[3]
_set_title(s, "DATASET")
_remove_content_ph(s)

# Stats text (top portion)
_add_textbox(s, CONTENT_L, CONTENT_T, CONTENT_W, 2.8, [
    {"text": "Sources",
     "bold": True, "size": 16, "color": BLUE_TEXT, "indent": 0},
    {"text": "Kermany Chest X-Ray (Kaggle)  —  5,863 images  |  Classes: Normal, Pneumonia",
     "bold": False, "size": 13, "indent": 1},
    {"text": "COVID-19 Radiography Database (Kaggle)  —  21,165 images  |  Classes: COVID, Viral Pneumonia, Lung Opacity, Normal",
     "bold": False, "size": 13, "indent": 1},

    {"text": "After hash-based deduplication \u2192 27,005 unique images",
     "bold": True, "size": 14, "indent": 0, "space_before": 6, "color": BLUE_TEXT},

    {"text": "80 / 10 / 10 stratified split   |   Train: 21,604   |   Val: 2,700   |   Test: 2,701",
     "bold": False, "size": 13, "indent": 1},
    {"text": "Normal: 11,767   |   Pneumonia: 11,622   |   COVID: 3,616  (\u226413.4% \u2014 minority class)",
     "bold": False, "size": 13, "indent": 1},

    {"text": "Preprocessing",
     "bold": True, "size": 16, "color": BLUE_TEXT, "indent": 0, "space_before": 6},
    {"text": "Resize to 224\u00d7224  |  ImageNet normalisation  |  Train augmentation: H-flip, \u00b115\u00b0 rotation, affine",
     "bold": False, "size": 13, "indent": 1},
])

# Sample X-ray images (bottom row)  —  ratios: 1.46, 1.51, 1.00
# Fixed height 2.4", widths from ratios, centered across 12.5"
img_h   = 2.4
widths  = [img_h * 1.46, img_h * 1.51, img_h * 1.00]   # 3.50, 3.62, 2.40
gap     = 0.25
total_w = sum(widths) + 2 * gap                          # 3.50+3.62+2.40+0.50 = 10.02
start_l = CONTENT_L + (CONTENT_W - total_w) / 2         # centre the row
img_t   = 4.85

for i, (fname, label) in enumerate([
    ("IM-0023-0001.jpeg",            "Normal"),
    ("person1946_bacteria_4875.jpeg","Pneumonia"),
    ("x-ray-image-2b_full.jpg",      "COVID-19"),
]):
    lx = start_l + sum(widths[:i]) + i * gap
    _add_img(s, SAMPLE / fname, lx, img_t, width=widths[i], height=img_h)
    _add_caption(s, label, lx, img_t + img_h + 0.03, widths[i], size=11)


# ── SLIDE 5 ── Experiments ────────────────────────────────────────────────────
s = slides[4]
_set_title(s, "EXPERIMENTS")
_remove_content_ph(s)

# Compact text block (top)
_add_textbox(s, CONTENT_L, CONTENT_T, CONTENT_W, 2.4, [
    {"text": "Hardware & Setup",
     "bold": True, "size": 16, "color": BLUE_TEXT, "indent": 0},
    {"text": "NVIDIA RTX 3060 Laptop (6 GB VRAM), CUDA 12.4  |  PyTorch + mixed precision (torch.cuda.amp)",
     "bold": False, "size": 13, "indent": 1},
    {"text": "Batch size: 32  |  Image size: 224\u00d7224  |  Optimizer: Adam",
     "bold": False, "size": 13, "indent": 1},

    {"text": "Phase Results",
     "bold": True, "size": 16, "color": BLUE_TEXT, "indent": 0, "space_before": 6},
    {"text": "Phase 1 final val-acc  \u2014  DenseNet121: 79.1%   ResNet50: 79.6%   EfficientNetB0: 86.6%",
     "bold": False, "size": 13, "indent": 1},
    {"text": "Phase 2 best  val-acc  \u2014  DenseNet121: 90.1%   ResNet50: 92.3%   EfficientNetB0: 94.1%",
     "bold": True,  "size": 13, "indent": 1},
    {"text": "EarlyStopping triggered: ResNet50 at Phase-1 ep 4 and Phase-2 ep 16; EfficientNetB0 at Phase-2 ep 7",
     "bold": False, "size": 12, "indent": 2},
])

# Training curves image — ratio 2.85; W=12.0 → H=4.21
tc_w, tc_h = 12.0, 12.0 / 2.85
_add_img(s, IMGS / "training_curves.png",
         left=CONTENT_L + (CONTENT_W - tc_w) / 2,
         top=4.05, width=tc_w, height=tc_h)
_add_caption(s, "Fig: Training & validation loss/accuracy across Phase 1 and Phase 2 for all three models",
             CONTENT_L, 4.05 + tc_h + 0.03, CONTENT_W, size=10)


# ── SLIDE 6 ── Results ────────────────────────────────────────────────────────
s = slides[5]
_set_title(s, "RESULTS WITH VISUAL OUTPUTS")
_remove_content_ph(s)

# ── Metrics table ──
rows_data = [
    # (Model, Acc, AUC, MacroF1, F1-COVID)
    ("Model",          "Accuracy", "Macro AUC", "Macro F1", "F1-COVID"),
    ("DenseNet121",    "90.19%",   "0.9814",    "0.8937",   "0.8678"),
    ("ResNet50",       "92.63%",   "0.9875",    "0.9232",   "0.9131"),
    ("EfficientNetB0", "93.67%",   "0.9900",    "0.9376",   "0.9407"),
    ("Ensemble",       "93.74%",   "0.9908",    "0.9396",   "0.9465"),
]
n_rows, n_cols = len(rows_data), len(rows_data[0])
tbl_w, tbl_h = CONTENT_W, 2.65
tf_obj = s.shapes.add_table(n_rows, n_cols,
                             Inches(CONTENT_L), Inches(CONTENT_T),
                             Inches(tbl_w), Inches(tbl_h))
tbl = tf_obj.table

# Column widths
col_widths = [3.0, 2.0, 2.0, 2.0, 2.0]   # total ≈ 11.0 → scaled to 12.5
scale = tbl_w / sum(col_widths)
for ci, cw in enumerate(col_widths):
    tbl.columns[ci].width = Inches(cw * scale)

for ri, row in enumerate(rows_data):
    is_header  = ri == 0
    is_best    = ri == 4   # ensemble
    bg = DARK_BLUE if is_header else (HIGHLIGHT if is_best else (LIGHT_BLUE if ri % 2 == 0 else WHITE))
    fg = WHITE_TEXT if is_header else DARK_TEXT
    for ci, val in enumerate(row):
        align = PP_ALIGN.CENTER if ci > 0 else PP_ALIGN.LEFT
        _set_cell(tbl, ri, ci, val,
                  bold=(is_header or is_best),
                  size=(12 if is_header else 12),
                  align=align,
                  color_rgb=fg, bg_hex=bg)

# Insight line
_add_textbox(s, CONTENT_L, CONTENT_T + tbl_h + 0.08, CONTENT_W, 0.35, [
    {"text": "High-confidence subset (55.2% of test set, threshold=0.35):  99.40% accuracy",
     "bold": True, "size": 13, "indent": 0, "color": (31, 78, 121)},
])

# Confusion matrices — ratio 4.39; W=12.0 → H=2.73
cm_w, cm_h = 12.0, 12.0 / 4.39
cm_t = CONTENT_T + tbl_h + 0.55
_add_img(s, IMGS / "confusion_matrices.png",
         left=CONTENT_L + (CONTENT_W - cm_w) / 2,
         top=cm_t, width=cm_w, height=cm_h)
_add_caption(s, "Fig: Normalised confusion matrices — DenseNet121, ResNet50, EfficientNetB0, Ensemble",
             CONTENT_L, cm_t + cm_h + 0.03, CONTENT_W, size=10)


# ── SLIDE 7 ── Comparison with SOTA ──────────────────────────────────────────
s = slides[6]
_set_title(s, "COMPARISON WITH STATE OF THE ART")
_remove_content_ph(s)

sota_rows = [
    # Method, Classes, Dataset, Acc, AUC, UQ / Explainability
    ("Method",              "Classes", "Dataset",  "Accuracy", "AUC",   "UQ / Explainability"),
    ("CheXNet (2017)",      "14-class","100K CXR",  "90.1%",  "0.968", "None"),
    ("COVID-Net (2020)",    "3-class", "~14K CXR",  "92.7%",  "—",     "None"),
    ("Minaee et al. 2021",  "2-class", "5K CXR",    "98.3%",  "—",     "Grad-CAM only"),
    ("Narin et al. 2021",   "3-class", "4.2K CXR",  "98.2%",  "0.993", "None"),
    ("Zhang et al. 2020",   "3-class", "1.4K CXR",  "96.0%",  "0.994", "Conf. only"),
    ("Ours (Ensemble)",     "3-class", "27K CXR",   "93.74%", "0.991", "MC Dropout + Grad-CAM"),
]
nr, nc = len(sota_rows), len(sota_rows[0])
st_w, st_h = CONTENT_W, 4.8
st_obj = s.shapes.add_table(nr, nc,
                              Inches(CONTENT_L), Inches(CONTENT_T),
                              Inches(st_w), Inches(st_h))
st = st_obj.table

col_w = [3.0, 1.5, 1.8, 1.6, 1.4, 2.8]
scale_s = st_w / sum(col_w)
for ci, cw in enumerate(col_w):
    st.columns[ci].width = Inches(cw * scale_s)

for ri, row in enumerate(sota_rows):
    is_header = ri == 0
    is_ours   = ri == 6
    bg = DARK_BLUE if is_header else (HIGHLIGHT if is_ours else (LIGHT_BLUE if ri % 2 == 0 else WHITE))
    fg = WHITE_TEXT if is_header else DARK_TEXT
    for ci, val in enumerate(row):
        align = PP_ALIGN.CENTER if ci > 0 else PP_ALIGN.LEFT
        _set_cell(st, ri, ci, val,
                  bold=(is_header or is_ours),
                  size=12,
                  align=align,
                  color_rgb=fg, bg_hex=bg)

# Key insight below table
_add_textbox(s, CONTENT_L, CONTENT_T + st_h + 0.12, CONTENT_W, 0.8, [
    {"text": "\u2605  Largest unified 3-class dataset (27K) among compared methods",
     "bold": False, "size": 12, "indent": 0, "color": DARK_TEXT},
    {"text": "\u2605  Only method combining uncertainty quantification AND visual explainability",
     "bold": False, "size": 12, "indent": 0, "color": DARK_TEXT},
    {"text": "\u2605  Referral mechanism achieves 99.40% accuracy on high-confidence predictions",
     "bold": True, "size": 12, "indent": 0, "color": BLUE_TEXT},
])


# ── SLIDE 8 ── Summary and Future Work ───────────────────────────────────────
s = slides[7]
_set_title(s, "SUMMARY AND FUTURE WORK")
_remove_content_ph(s)

_add_textbox(s, CONTENT_L, CONTENT_T, 6.0, CONTENT_H, [
    {"text": "Summary", "bold": True, "size": 17, "color": BLUE_TEXT, "indent": 0},
    {"text": "End-to-end 3-class CXR classifier on 27,005 images",
     "bold": False, "size": 13, "indent": 1},
    {"text": "Ensemble (DenseNet121 + ResNet50 + EfficientNetB0)",
     "bold": False, "size": 13, "indent": 1},
    {"text": "93.74% accuracy  |  0.9908 macro AUC  |  0.9396 macro F1",
     "bold": True, "size": 13, "indent": 2},
    {"text": "Focal Loss + class weights handle severe class imbalance",
     "bold": False, "size": 13, "indent": 1},
    {"text": "MC Dropout: 99.40% accuracy on 55.2% high-confidence subset",
     "bold": False, "size": 13, "indent": 1},
    {"text": "Grad-CAM provides radiologist-interpretable saliency overlays",
     "bold": False, "size": 13, "indent": 1},
    {"text": "Deployed as Streamlit app: single & batch inference + CSV export",
     "bold": False, "size": 13, "indent": 1},

    {"text": "Limitations", "bold": True, "size": 17, "color": BLUE_TEXT,
     "indent": 0, "space_before": 8},
    {"text": "Two-source dataset — possible domain gap between Kermany and COVID-19 DB",
     "bold": False, "size": 13, "indent": 1},
    {"text": "44.8% referral rate may be high for real-world triage deployment",
     "bold": False, "size": 13, "indent": 1},

    {"text": "Future Work", "bold": True, "size": 17, "color": BLUE_TEXT,
     "indent": 0, "space_before": 8},
    {"text": "Self-supervised / foundation model pretraining on unlabelled CXRs",
     "bold": False, "size": 13, "indent": 1},
    {"text": "Vision-language model for automated radiology report generation",
     "bold": False, "size": 13, "indent": 1},
    {"text": "Reduce referral rate via temperature scaling / post-hoc calibration",
     "bold": False, "size": 13, "indent": 1},
    {"text": "Prospective clinical validation and DICOM integration",
     "bold": False, "size": 13, "indent": 1},
])

# ROC curves on right — ratio 3.63; fix H=5.5 → W=19.97 too wide
# Use W=6.0 → H=1.65, not ideal. Better: W=6.0, H=3.5 (slight stretch)
roc_side_h = 2.8
roc_side_w = roc_side_h * 3.63   # 10.16" — too wide for right panel
# Clamp to 6.2" and allow slight distortion
roc_side_w = 6.0
_add_img(s, IMGS / "roc_curves.png", 6.6, CONTENT_T, width=roc_side_w, height=roc_side_h)
_add_caption(s, "ROC Curves (per class, all models)",
             6.6, CONTENT_T + roc_side_h + 0.03, roc_side_w, size=10)

unc_h = 2.6
unc_w = 6.0
_add_img(s, IMGS / "uncertainty_analysis.png",
         6.6, CONTENT_T + roc_side_h + 0.35,
         width=unc_w, height=unc_h)
_add_caption(s, "Uncertainty Analysis (MC Dropout, threshold=0.35)",
             6.6, CONTENT_T + roc_side_h + 0.35 + unc_h + 0.03,
             unc_w, size=10)


# ── SLIDE 9 ── Visual Analysis ────────────────────────────────────────────────
s = slides[8]
_set_title(s, "VISUAL ANALYSIS — CALIBRATION & EXPLAINABILITY")
_remove_content_ph(s)

# 2×2 image grid
CELL_W = 6.0
CELL_H = 2.45
GAP    = 0.35
ROW1_T = CONTENT_T
ROW2_T = ROW1_T + CELL_H + GAP      # 1.55 + 2.45 + 0.35 = 4.35
COL1_L = CONTENT_L                  # 0.40
COL2_L = COL1_L + CELL_W + GAP     # 0.40 + 6.0 + 0.35 = 6.75

# Row 1: training curves | calibration diagrams
_add_img(s, IMGS / "training_curves.png",    COL1_L, ROW1_T, CELL_W, CELL_H)
_add_caption(s, "Training & Validation Curves (Loss / Accuracy)",
             COL1_L, ROW1_T + CELL_H + 0.02, CELL_W)

_add_img(s, IMGS / "calibration_diagrams.png", COL2_L, ROW1_T, CELL_W, CELL_H)
_add_caption(s, "Calibration (Reliability) Diagrams — all models & ensemble",
             COL2_L, ROW1_T + CELL_H + 0.02, CELL_W)

# Row 2: ROC curves | uncertainty analysis
_add_img(s, IMGS / "roc_curves.png",         COL1_L, ROW2_T, CELL_W, CELL_H)
_add_caption(s, "One-vs-Rest ROC Curves (per class, all models)",
             COL1_L, ROW2_T + CELL_H + 0.02, CELL_W)

_add_img(s, IMGS / "uncertainty_analysis.png", COL2_L, ROW2_T, CELL_W, CELL_H)
_add_caption(s, "MC Dropout Uncertainty — correct vs misclassified (threshold=0.35)",
             COL2_L, ROW2_T + CELL_H + 0.02, CELL_W)


# ═══════════════════════════════════════════════════════════════════════════════
prs.save(str(OUT))
print(f"Saved -> {OUT}")
print(f"Slides: {len(prs.slides)}")
